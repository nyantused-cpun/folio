# -*- coding: utf-8 -*-
"""
文件读取桥接层。
MD/TXT 由 AI 原生 Read 工具直接读取。
PDF/DOCX/PPTX 通过本脚本转换为纯文本后交 AI 理解。
不需要任何重型解析引擎（如 Docling）。
"""
import os
import sys

SCRIPT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # 仓库根


def read_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def read_docx(path):
    from docx import Document
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text)
    for table in doc.tables:
        lines.append("--- 表格 ---")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            lines.append(" | ".join(cells))
    return "\n\n".join(lines)


def read_pdf(path):
    # P1-4：用 try/finally 确保 doc.close() 在异常时也执行
    import fitz
    doc = fitz.open(path)
    try:
        lines = []
        for i, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                lines.append(f"\n=== 第 {i} 页 ===")
                lines.append(text)
        return "\n".join(lines) if lines else "(PDF 无文本内容)"
    finally:
        doc.close()


def read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append(" | ".join(c.text.strip() for c in row.cells))
                slide_text.append("表格: " + " / ".join(rows))
        if slide_text:
            lines.append(f"\n=== 第 {i} 页 ===")
            lines.append("\n".join(slide_text))
    return "\n".join(lines)


def read_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True, read_only=True)
    try:
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"\n=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    finally:
        wb.close()


def read_xlsx_structured(path):
    """xlsx 结构化读取：每行作为一个 chunk，保留 sheet 名和行号。

    返回 [(chunk_text, anchor), ...]
      - chunk_text: "Sheet名 | 字段1 | 字段2 | ..."（带表头上下文）
      - anchor: "sheet{N}#row{M}"

    相比 read_xlsx 的纯文本输出，此格式让报价单的每一行保持完整，
    价格、产品名、规格不会被滑动窗口切断。
    """
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True, read_only=True)
    try:
        chunks = []
        for s_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            header = None
            for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
                cells = [str(c) if c is not None else "" for c in row]
                if not any(cells):
                    continue
                # 第一行非空作为表头（报价单常见结构）
                if header is None and r_idx == 0:
                    header = cells
                # 拼接：Sheet名 + 表头 + 当前行，保留上下文
                parts = [f"[Sheet: {sheet_name}]"]
                if header and r_idx > 0:
                    # 把表头和当前行配对，只保留非空字段
                    pairs = [(h, v) for h, v in zip(header, cells, strict=False) if v]
                    if pairs:
                        parts.append(" | ".join(f"{h}={v}" for h, v in pairs))
                    else:
                        parts.append(" | ".join(cells))
                else:
                    parts.append(" | ".join(cells))
                chunk_text = " ".join(parts)
                chunks.append((chunk_text, f"sheet{s_idx}#row{r_idx}"))
        return chunks
    finally:
        wb.close()


def read_image(path):
    """图片 → 文本描述（调云端视觉模型）。"""
    from _cloud_llm import vision_chat
    prompt = "详细描述这张图片的内容，包括：图表数据、流程结构、关键文字、UI 元素。如果是数据图表，把数据表格化输出。"
    desc = vision_chat(prompt, path)
    return desc if desc else "(图片视觉理解失败，需 AI 直接查看)"


READERS = {
    ".txt": read_txt, ".md": read_txt,
    ".docx": read_docx,
    ".pdf": read_pdf,
    ".pptx": read_pptx,
    ".xlsx": read_xlsx,
    ".jpg": read_image, ".jpeg": read_image, ".png": read_image,
    ".gif": read_image, ".bmp": read_image, ".webp": read_image,
    ".svg": read_image, ".tiff": read_image, ".tif": read_image, ".heic": read_image,
}


import hashlib

CACHE_DIR = os.path.join(SCRIPT_DIR, "_knowledge", ".cache")


def read_full(path):
    """全文读取，写入缓存，返回 (摘要, 缓存路径)。
    不截断，AI 可按需 Read 缓存文件对应 offset 读全文。"""
    if not os.path.isabs(path):
        path = os.path.join(SCRIPT_DIR, path)

    ext = os.path.splitext(path)[1].lower()
    reader = READERS.get(ext)
    if not reader:
        supported = ", ".join(sorted(READERS.keys()))
        return f"不支持的格式: {ext}（支持: {supported}）", ""

    text = reader(path)

    mtime = os.path.getmtime(path)
    file_hash = hashlib.md5(f"{path}{mtime}".encode("utf-8"), usedforsecurity=False).hexdigest()[:16]
    os.makedirs(CACHE_DIR, exist_ok=True)
    cache_path = os.path.join(CACHE_DIR, f"{file_hash}.txt")

    with open(cache_path, "w", encoding="utf-8") as f:
        f.write(text)

    try:
        from _compress import summarize_for_attention
        summary = summarize_for_attention(text)
    except ImportError:
        summary = text[:2000] + (f"\n... (共 {len(text)} 字符)" if len(text) > 2000 else "")

    return summary, cache_path


def read_batch(paths, max_workers=4):
    """并行批量读取多个文件，返回 [(path, summary, cache), ...]。"""
    from concurrent.futures import ThreadPoolExecutor, as_completed
    results = [None] * len(paths)

    def _worker(idx, p):
        try:
            s, c = read_full(p)
            return idx, (p, s, c)
        except Exception as e:
            return idx, (p, f"读取失败: {e}", "")

    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = [ex.submit(_worker, i, p) for i, p in enumerate(paths)]
        for fut in as_completed(futures):
            idx, data = fut.result()
            results[idx] = data
    return results


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("文件文本提取（供 AI 理解消费）")
        print("用法:")
        print("  python _pipeline.py <文件路径>              原始提取（截断 15000）")
        print("  python _pipeline.py --full <文件路径>       全文缓存模式（返回摘要+缓存路径）")
        print("支持: .txt .md .docx .pdf .pptx .xlsx")
        sys.exit(0)

    if sys.argv[1] == "--full" and len(sys.argv) >= 3:
        path = sys.argv[2]
        summary, cache = read_full(path)
        print("=== 注意力摘要 ===")
        print(summary)
        print(f"\n=== 缓存路径 ===\n{cache}")
        print(f"\n=== 提示 ===\nAI 可 Read 此缓存读全文：{cache}")
        sys.exit(0)

    path = sys.argv[1]
    if not os.path.isabs(path):
        path = os.path.join(SCRIPT_DIR, path)

    if not os.path.exists(path):
        print(f"文件不存在: {path}")
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()
    reader = READERS.get(ext)

    if not reader:
        print(f"不支持的格式: {ext}")
        sys.exit(1)

    text = reader(path)
    if len(text) > 15000:
        print(text[:15000])
        print(f"\n... (截断，原文共 {len(text)} 字符)")
    else:
        print(text)
