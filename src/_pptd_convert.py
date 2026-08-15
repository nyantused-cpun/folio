# -*- coding: utf-8 -*-
"""pptd v1 方言 -> PPTX 转换器（python-pptx 实现，C 批次自研路线）。

背景：kimi_ppt_dsl.pyz 已删、kimi-slides.exe 本机 CLI 不可用（C-1 实测），
2026-08-12 决策改用 python-pptx 自研转换后端，保留 v1 pptd 方言与
_pptd_gen.py 生成器零改动。

本模块提供 backend 的两个能力：
- check_pptd(pptd_path) -> (ok, warn_counts)  本地结构校验（error 阻断、
  warning 放行；深度版式检查由 _verify.verify_pptd_deck / _layout_lint 负责）
- convert_pptd(pptd_path, out_path) -> bool   工程 -> PPTX

v1 方言要点（迁移基准见 docs/材料生成管线C1_方言迁移差异清单_2026-08-12.md）：
- 主文件：title / size([1280,720]) / theme / pages（无 version 字段）
- 页面：pageType / background / elements[]
- 元素：elementId / elementType / bounds + 类型字段
  - text: content{style|显式样式, align, wrap, text(富文本 <p>/<strong>/<em>)}
  - shape: shapeName + adjustments + fill/border + flip
  - 连接线: elementType=shape + shapeName(straightConnector1/bentConnector3/
    curvedConnector3) + bounds(端点包围盒) + flip + arrow[start,end]
  - table: columnWidths / rowHeights / style / rows[[{content:{...}}]]
  - image: src 绝对路径 + fit{mode}
"""
import os

import yaml

# ---------------------------------------------------------------------------
# 常量与映射
# ---------------------------------------------------------------------------
EMU_PER_PT = 12700  # python-pptx 单位：1pt = 12700 EMU；v1 方言 1px = 1pt

# 转换器支持的元素类型（chart 不在生成器范围，C-3 不支持，check 已拦）
_SUPPORTED_TYPES = ("text", "shape", "table", "image", "group")

# v1 shapeName -> python-pptx MSO_SHAPE 成员名（运行时 getattr 解析，防版本漂移）。
# 只收录 v1 生成器实际产出的 preset + D-2 明确会引入的块箭头/语义形状；
# flowChart* 家族（D-2 落地，语义化 preset：terminator/process/decision/document）。
_SHAPE_NAME_MAP = {
    # v1 实际产出（grep 核实）
    "rect": "RECTANGLE",
    "roundRect": "ROUNDED_RECTANGLE",
    "snip2SameRect": "SNIP_2_SAME_RECTANGLE",
    "ellipse": "OVAL",
    "triangle": "ISOSCELES_TRIANGLE",
    "diamond": "DIAMOND",
    "chevron": "CHEVRON",
    "rightArrow": "RIGHT_ARROW",
    # D-2 明确引入（块箭头扩充 / 语义 preset）
    "hexagon": "HEXAGON",
    "pentagon": "REGULAR_PENTAGON",
    "star5": "STAR_5_POINT",
    "arc": "ARC",
    "leftArrow": "LEFT_ARROW",
    "upArrow": "UP_ARROW",
    "downArrow": "DOWN_ARROW",
    "leftRightArrow": "LEFT_RIGHT_ARROW",
    "upDownArrow": "UP_DOWN_ARROW",
    "parallelogram": "PARALLELOGRAM",
    "trapezoid": "TRAPEZOID",
    "pie": "PIE",
    "blockArc": "BLOCK_ARC",
    "foldedCorner": "FOLDED_CORNER",
    # data_flow 数据库罐形（relationship.py）
    "can": "CAN",
    # D-2 语义化 preset（flow 节点：start/end/process/decision/doc）
    "flowChartTerminator": "FLOWCHART_TERMINATOR",
    "flowChartProcess": "FLOWCHART_PROCESS",
    "flowChartDecision": "FLOWCHART_DECISION",
    "flowChartDocument": "FLOWCHART_DOCUMENT",
}

# 连接线 shapeName -> MSO_CONNECTOR 成员名（v1 connector 家族）
_CONNECTOR_MAP = {
    "straightConnector1": "STRAIGHT",
    "bentConnector3": "ELBOW",
    "curvedConnector3": "CURVE",
}

# v1 箭头名 -> OOXML lineEnd type（a:headEnd/a:tailEnd 的 type 枚举，ST_LineEndType）
_ARROW_MAP = {
    "none": None,
    "arrow": "triangle",
    "stealth": "stealth",
    "diamond": "diamond",
    "oval": "oval",
}

# 水平对齐映射
_ALIGN_H_MAP = {"left": "LEFT", "center": "CENTER", "right": "RIGHT",
                "justify": "JUSTIFY"}
# 垂直对齐映射
_ANCHOR_V_MAP = {"top": "TOP", "middle": "MIDDLE", "bottom": "BOTTOM"}


# ---------------------------------------------------------------------------
# YAML 读取
# ---------------------------------------------------------------------------
def load_yaml_file(path):
    """读取 YAML 文件；解析失败抛异常（调用方转为 check error）。"""
    with open(path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    return data if isinstance(data, dict) else {}


# ---------------------------------------------------------------------------
# check：本地结构校验（error 阻断，warning 放行）
# ---------------------------------------------------------------------------
def check_pptd(pptd_path):
    """基础结构校验：主文件可解析、必填字段、页面文件存在、元素结构合法。

    返回 (ok, warn_counts)。五类 Warning 计数恒为空 dict——版式维检查
    （TextOverflow 等）由 verify_pptd_deck / lint_pptd_dir 在 build 链路
    承担，这里只做结构门。
    """
    warn_counts = {}
    deck_dir = os.path.dirname(os.path.abspath(pptd_path)) or "."

    # 1. 主文件可解析 + 必填字段
    try:
        main = load_yaml_file(pptd_path)
    except Exception as e:
        print(f"[check error] 主文件解析失败: {e}")
        return False, warn_counts
    for field in ("size", "pages"):
        if field not in main:
            print(f"[check error] 主文件缺少必填字段: {field}")
            return False, warn_counts
    size = main.get("size")
    if not (isinstance(size, (list, tuple)) and len(size) == 2):
        print(f"[check error] size 必须是 [宽, 高] 二元数组: {size!r}")
        return False, warn_counts

    # 2. 页面文件存在 + 每页元素结构
    pages = main.get("pages") or []
    for rel in pages:
        page_path = os.path.join(deck_dir, rel)
        if not os.path.exists(page_path):
            print(f"[check error] 页面文件不存在: {rel}")
            return False, warn_counts
        try:
            page = load_yaml_file(page_path)
        except Exception as e:
            print(f"[check error] 页面解析失败 {rel}: {e}")
            return False, warn_counts
        elements = page.get("elements")
        if not isinstance(elements, list):
            print(f"[check error] 页面缺少 elements 数组: {rel}")
            return False, warn_counts
        seen_ids = set()
        for i, el in enumerate(elements):
            if not isinstance(el, dict):
                print(f"[check error] {rel} elements[{i}] 不是对象")
                return False, warn_counts
            for field in ("elementId", "elementType", "bounds"):
                if field not in el:
                    print(f"[check error] {rel} elements[{i}] 缺少必填字段: {field}")
                    return False, warn_counts
            eid = el["elementId"]
            if eid in seen_ids:
                print(f"[check error] {rel} elementId 重复: {eid}")
                return False, warn_counts
            seen_ids.add(eid)
            bounds = el["bounds"]
            if not (isinstance(bounds, (list, tuple)) and len(bounds) == 4):
                print(f"[check error] {rel} {eid} bounds 必须是 [x,y,w,h]: {bounds!r}")
                return False, warn_counts
            if el["elementType"] not in _SUPPORTED_TYPES:
                print(f"[check error] {rel} {eid} 未知 elementType: "
                      f"{el['elementType']}（支持: {'/'.join(_SUPPORTED_TYPES)}）")
                return False, warn_counts

    print("[check] 基础结构校验通过")
    return True, warn_counts


# ---------------------------------------------------------------------------
# convert：工程 -> PPTX（python-pptx）
# ---------------------------------------------------------------------------
def convert_pptd(pptd_path, out_path):
    """把 v1 pptd 工程转换为 PPTX（python-pptx）。

    覆盖元素：text / shape（含 connector 连接线）/ table / image + 页面背景。
    主题（theme.colors/textStyles/tableStyles）解析到元素级样式。
    """
    from pptx import Presentation
    from pptx.util import Emu

    deck_dir = os.path.dirname(os.path.abspath(pptd_path)) or "."
    main = load_yaml_file(pptd_path)

    size = main.get("size") or [1280, 720]
    theme = main.get("theme") or {}
    theme_colors = theme.get("colors") or {}
    text_styles = theme.get("textStyles") or {}
    table_styles = theme.get("tableStyles") or {}

    prs = Presentation()
    prs.slide_width = Emu(int(size[0]) * EMU_PER_PT)
    prs.slide_height = Emu(int(size[1]) * EMU_PER_PT)

    blank_layout = prs.slide_layouts[6]  # Blank
    pages = main.get("pages") or []
    for idx, rel in enumerate(pages):
        page_path = os.path.join(deck_dir, rel)
        if not os.path.exists(page_path):
            print(f"[convert warning] 页面文件不存在，跳过: {rel}")
            continue
        page = load_yaml_file(page_path)
        slide = prs.slides.add_slide(blank_layout)
        _render_background(slide, page.get("background"), theme_colors, prs)
        for el in page.get("elements") or []:
            try:
                _render_element(slide, el, theme_colors, text_styles,
                                table_styles, deck_dir)
            except Exception as e:
                print(f"[convert warning] {rel} 元素 {el.get('elementId')} "
                      f"渲染失败，跳过: {e}")
        print(f"[convert] 第 {idx + 1}/{len(pages)} 页已渲染: {rel}")

    os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
    prs.save(out_path)
    print(f"[convert] 已生成: {out_path}")
    return True


def _render_background(slide, bg, theme_colors, prs):
    """页面背景：solid（默认白）或 image（铺满）。"""
    if not isinstance(bg, dict):
        bg = {"type": "solid", "color": "#FFFFFF"}
    btype = bg.get("type", "solid")
    if btype == "solid":
        color = _color_value(bg.get("color", "#FFFFFF"), theme_colors)
        if color:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = color
    elif btype == "image":
        src = bg.get("src")
        if src and os.path.exists(src):
            from pptx.util import Emu
            slide.shapes.add_picture(
                src, Emu(0), Emu(0), prs.slide_width, prs.slide_height)


def _render_element(slide, el, theme_colors, text_styles, table_styles, deck_dir):
    """按 elementType 分发渲染。返回渲染出的 shape 对象（group 组合用）。"""
    etype = el["elementType"]
    if etype == "text":
        return _render_text(slide, el, theme_colors, text_styles)
    elif etype == "shape":
        if el.get("shapeName") in _CONNECTOR_MAP:
            return _render_connector(slide, el, theme_colors)
        else:
            return _render_shape(slide, el, theme_colors)
    elif etype == "table":
        return _render_table(slide, el, theme_colors, table_styles)
    elif etype == "image":
        return _render_image(slide, el, deck_dir)
    elif etype == "group":
        return _render_group(slide, el, theme_colors, text_styles, table_styles, deck_dir)
    return None


def _render_group(slide, el, theme_colors, text_styles, table_styles, deck_dir):
    """组合：先渲染 children 到 slide 收集 shape 对象，再 add_group_shape 组合。

    角标/框体/色条/标题等节点内部元素用 group 组合成一个整体，PPT 里
    拖动/缩放时保持相对关系不变（python-pptx 自动 recalculate 包围盒）。
    """
    shapes = []
    for child in el.get("children") or []:
        obj = _render_element(slide, child, theme_colors, text_styles, table_styles, deck_dir)
        if obj is not None:
            shapes.append(obj)
    if shapes:
        slide.shapes.add_group_shape(shapes)
    return None


def _color_value(value, theme_colors):
    """解析颜色值为 RGBColor；$token 从 theme 取；无效返回 None。"""
    from pptx.dml.color import RGBColor
    if not value:
        return None
    s = str(value).strip()
    if s.startswith("$"):
        s = theme_colors.get(s[1:])
        if not s:
            return None
        s = str(s).strip()
    if s.startswith("#"):
        s = s[1:]
    if len(s) == 6 and all(c in "0123456789abcdefABCDEF" for c in s):
        return RGBColor.from_string(s)
    # 8 位 hex（带 alpha）：取前 6 位
    if len(s) == 8 and all(c in "0123456789abcdefABCDEF" for c in s):
        return RGBColor.from_string(s[:6])
    return None


def _style_dict(style_key, text_styles):
    """取主题 textStyle（$key），返回 dict 或 None。"""
    if style_key and isinstance(style_key, str) and style_key.startswith("$"):
        return text_styles.get(style_key[1:]) or {}
    return {}


# ---------------------------------------------------------------------------
# text 元素
# ---------------------------------------------------------------------------
def _render_text(slide, el, theme_colors, text_styles):
    """文本框：content.style（主题样式）+ 显式字段覆盖 + 富文本 <p>/<strong>/<em>。"""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    x, y, w, h = [int(v) for v in el["bounds"]]
    content = el.get("content") or {}
    tb = slide.shapes.add_textbox(Emu(x * EMU_PER_PT), Emu(y * EMU_PER_PT),
                                  Emu(w * EMU_PER_PT), Emu(h * EMU_PER_PT))
    tf = tb.text_frame
    tf.word_wrap = bool(content.get("wrap", True))
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # 样式基线：主题引用 + 显式字段覆盖
    base = dict(_style_dict(content.get("style"), text_styles))
    for k in ("fontSize", "color", "fontFamily", "lineHeight"):
        if k in content:
            base[k] = content[k]

    # 垂直对齐
    align = content.get("align") or ["left", "top"]
    if isinstance(align, (list, tuple)) and len(align) >= 2:
        anchor = _ANCHOR_V_MAP.get(str(align[1]).lower())
        if anchor:
            tf.vertical_anchor = getattr(MSO_ANCHOR, anchor)

    paragraphs = _parse_rich_text(content.get("text") or "")
    if not paragraphs:
        paragraphs = [""]
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        halign = _ALIGN_H_MAP.get(str(align[0]).lower())
        if halign:
            p.alignment = getattr(PP_ALIGN, halign)
        if base.get("lineHeight"):
            p.line_spacing = float(base["lineHeight"])
        for text, bold, italic in (runs or [("", False, False)]):
            run = p.add_run()
            run.text = text
            run.font.size = Pt(int(base.get("fontSize", 18)))
            run.font.bold = bold or bool(base.get("bold"))
            run.font.italic = italic or bool(base.get("italic"))
            color = _color_value(base.get("color"), theme_colors)
            if color:
                run.font.color.rgb = color
            family = base.get("fontFamily")
            if family:
                run.font.name = family
    return tb


def _unescape_html(text):
    """还原 _esc 转义的 HTML 实体（PPT 端不解析实体，需还原为原字符）。"""
    return (str(text).replace("&quot;", '"').replace("&lt;", "<")
            .replace("&gt;", ">").replace("&amp;", "&"))


def _parse_rich_text(text):
    """解析 v1 富文本 -> 段落列表，每段是 [(run_text, bold, italic)]。

    支持标签：<p>（分段）、<strong>/<b>（粗体）、<em>/<i>（斜体）。
    v1 的 text 字段 body 固定末尾带 ``\\n``（pe.text 拼接），无 <p> 时整段
    strip 掉该后缀与 YAML 缩进残留；分段后每 run 也 strip 首尾空白。
    """
    import re
    if not text:
        return []
    # 以 <p> 分段；无 <p> 时整段处理
    text = text.replace("\r", "")
    if "<p" not in text:
        paragraphs = [text.strip()]
    else:
        paragraphs = re.split(r"</?p[^>]*>", text)
        paragraphs = [p for p in paragraphs if p.strip()]

    result = []
    for para in paragraphs:
        if not para.strip():
            continue
        runs = []
        # 用正则逐段切 run：<strong>/<em> 包裹或纯文本
        pattern = re.compile(r"(<strong>|<b>|<em>|<i>|</strong>|</b>|</em>|</i>)")
        bold = italic = False
        buf = ""
        tokens = pattern.split(para)
        for tok in tokens:
            if tok == "<strong>" or tok == "<b>":
                if buf:
                    runs.append((buf, bold, italic))
                    buf = ""
                bold = True
            elif tok == "</strong>" or tok == "</b>":
                if buf:
                    runs.append((buf, bold, italic))
                    buf = ""
                bold = False
            elif tok == "<em>" or tok == "<i>":
                if buf:
                    runs.append((buf, bold, italic))
                    buf = ""
                italic = True
            elif tok == "</em>" or tok == "</i>":
                if buf:
                    runs.append((buf, bold, italic))
                    buf = ""
                italic = False
            else:
                buf += tok
        if buf:
            runs.append((buf, bold, italic))
        result.append([(_unescape_html(t.strip()), b, i) for t, b, i in runs if t.strip()]
                      or [("", False, False)])
    return result


# ---------------------------------------------------------------------------
# shape 元素
# ---------------------------------------------------------------------------
def _render_shape(slide, el, theme_colors):
    """形状：shapeName -> MSO_SHAPE preset；fill/border/adjustments。"""
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE

    x, y, w, h = [int(v) for v in el["bounds"]]
    name = el.get("shapeName") or "rect"
    member = _SHAPE_NAME_MAP.get(name)
    if not member or not hasattr(MSO_SHAPE, member):
        print(f"[convert warning] 未知 shapeName {name!r}，按矩形近似")
        member = "RECTANGLE"
    shape = slide.shapes.add_shape(
        getattr(MSO_SHAPE, member), Emu(x * EMU_PER_PT), Emu(y * EMU_PER_PT),
        Emu(w * EMU_PER_PT), Emu(h * EMU_PER_PT))

    # fill
    fill = el.get("fill")
    if isinstance(fill, dict) and fill.get("type") == "solid":
        color = _color_value(fill.get("color"), theme_colors)
        if color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
        else:
            shape.fill.background()
    else:
        shape.fill.background()

    # border
    border = el.get("border")
    if isinstance(border, dict):
        color = _color_value(border.get("color"), theme_colors)
        if color:
            shape.line.color.rgb = color
        else:
            shape.line.fill.background()
        if border.get("width"):
            shape.line.width = Pt(float(border["width"]))
        if border.get("style") == "dash":
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else:
        shape.line.fill.background()

    # adjustments（如 roundRect 圆角，v1 存 OOXML 原始值 0~100000，python-pptx 归一化 0~1）
    adjustments = el.get("adjustments")
    if adjustments and len(shape.adjustments) > 0:
        try:
            for i, val in enumerate(adjustments):
                if i < len(shape.adjustments):
                    shape.adjustments[i] = float(val) / 100000.0
        except (ValueError, IndexError):
            pass

    # flip（v1 存 [flipH, flipV]，如 timeline 红三角 flipV）
    flip = el.get("flip")
    if flip:
        _apply_flip(shape, bool(flip[0]), bool(flip[1]))

    # 节点框体（flowChart 系列）加柔和阴影提升层次；其余形状关闭继承阴影
    if name.startswith("flowChart"):
        _apply_soft_shadow(shape)
    else:
        shape.shadow.inherit = False
    return shape


def _apply_soft_shadow(shape, color="0F172A", opacity=32000, blur=25400,
                       dist=25400, direction=5400000):
    """给 shape 加柔和外阴影（操作 a:effectLst/a:outerShdw XML）。

    python-pptx 的 ShadowFormat 仅 expose inherit，无法设颜色/偏移/模糊，
    故直接写 OOXML。opacity 0~100000（20000=12%），blur/dist 单位 EMU。
    """
    from pptx.oxml.ns import qn
    sp_pr = shape._element.spPr
    effect_lst = sp_pr.find(qn("a:effectLst"))
    if effect_lst is None:
        effect_lst = sp_pr.makeelement(qn("a:effectLst"), {})
        sp_pr.append(effect_lst)
    for shdw in effect_lst.findall(qn("a:outerShdw")):
        effect_lst.remove(shdw)
    outer = effect_lst.makeelement(qn("a:outerShdw"), {
        "blurRad": str(blur), "dist": str(dist), "dir": str(direction),
        "rotWithShape": "0"})
    clr = outer.makeelement(qn("a:srgbClr"), {"val": color})
    alpha = outer.makeelement(qn("a:alpha"), {"val": str(opacity)})
    clr.append(alpha)
    outer.append(clr)
    effect_lst.append(outer)


def _apply_flip(shape, flip_h, flip_v):
    """设置形状水平/垂直翻转（python-pptx 无 flip API，操作 a:xfrm 的 flipH/flipV）。"""
    from pptx.oxml.ns import qn
    if not flip_h and not flip_v:
        return
    sp_pr = shape._element.spPr
    xfrm = sp_pr.find(qn("a:xfrm"))
    if xfrm is None:
        xfrm = sp_pr.makeelement(qn("a:xfrm"), {})
        sp_pr.insert(0, xfrm)
    if flip_h:
        xfrm.set("flipH", "1")
    if flip_v:
        xfrm.set("flipV", "1")


# ---------------------------------------------------------------------------
# connector 连接线（v1: elementType=shape + connector 家族 shapeName）
# ---------------------------------------------------------------------------
def _render_connector(slide, el, theme_colors):
    """连接线：bounds 是两端点包围盒，flip 表达方向，arrow 表达端点箭头。

    v1 契约（pptd_emit.connector）：bounds=[min_x,min_y,bw,bh]；
    flip_h=x2<x1、flip_v=y2<y1，据此还原起终点，保证箭头指向正确目标。
    """
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_CONNECTOR

    x, y, w, h = [int(v) for v in el["bounds"]]
    name = el.get("shapeName") or "straightConnector1"
    member = _CONNECTOR_MAP.get(name, "STRAIGHT")

    flip_h = bool((el.get("flip") or [False, False])[0])
    flip_v = bool((el.get("flip") or [False, False])[1])
    start_x = x + w if flip_h else x
    start_y = y + h if flip_v else y
    end_x = x if flip_h else x + w
    end_y = y if flip_v else y + h

    conn = slide.shapes.add_connector(
        getattr(MSO_CONNECTOR, member),
        Emu(start_x * EMU_PER_PT), Emu(start_y * EMU_PER_PT),
        Emu(end_x * EMU_PER_PT), Emu(end_y * EMU_PER_PT))

    border = el.get("border")
    if isinstance(border, dict):
        color = _color_value(border.get("color"), theme_colors)
        if color:
            conn.line.color.rgb = color
        if border.get("width"):
            conn.line.width = Pt(float(border["width"]))
        if border.get("style") == "dash":
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # 箭头（headEnd/tailEnd XML）
    arrow = el.get("arrow") or ["none", "none"]
    _set_arrow(conn, arrow)
    return conn


def _set_arrow(conn, arrow):
    """设置连接线端点箭头。arrow = [起点, 终点]，值见 _ARROW_MAP。"""
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    start_type = _ARROW_MAP.get(str(arrow[0]), None)
    end_type = _ARROW_MAP.get(str(arrow[1]), None)
    if start_type:
        head = ln.find(qn("a:headEnd"))
        if head is None:
            head = ln.makeelement(qn("a:headEnd"), {})
            ln.append(head)
        head.set("type", start_type)
    if end_type:
        tail = ln.find(qn("a:tailEnd"))
        if tail is None:
            tail = ln.makeelement(qn("a:tailEnd"), {})
            ln.append(tail)
        tail.set("type", end_type)


# ---------------------------------------------------------------------------
# table 元素
# ---------------------------------------------------------------------------
def _render_table(slide, el, theme_colors, table_styles):
    """表格：columnWidths/rowHeights 比例布局 + style($default) 主题样式。"""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN

    x, y, w, h = [int(v) for v in el["bounds"]]
    rows = el.get("rows") or []
    if not rows:
        return
    n_rows = len(rows)
    n_cols = len(rows[0])

    gframe = slide.shapes.add_table(
        n_rows, n_cols, Emu(x * EMU_PER_PT), Emu(y * EMU_PER_PT),
        Emu(w * EMU_PER_PT), Emu(h * EMU_PER_PT))
    table = gframe.table

    col_widths = el.get("columnWidths") or [1.0 / n_cols] * n_cols
    row_heights = el.get("rowHeights") or [1.0 / n_rows] * n_rows
    for i, ratio in enumerate(col_widths):
        if i < n_cols:
            table.columns[i].width = Emu(int(float(ratio) * w * EMU_PER_PT))
    for i, ratio in enumerate(row_heights):
        if i < n_rows:
            table.rows[i].height = Emu(int(float(ratio) * h * EMU_PER_PT))

    # 主题表样式（$default 引用）
    style_key = el.get("style")
    ts = {}
    if style_key and isinstance(style_key, str) and style_key.startswith("$"):
        ts = table_styles.get(style_key[1:]) or {}
    font_size = ts.get("fontSize", 14)
    font_family = ts.get("fontFamily", "Microsoft YaHei")
    header_fill = _color_value(ts.get("headerFill"), theme_colors)
    header_color = _color_value(ts.get("headerColor"), theme_colors)
    header_bold = bool(ts.get("headerBold", True))
    body_color = _color_value(ts.get("bodyColor"), theme_colors)
    border = ts.get("border")

    for r in range(n_rows):
        row = rows[r]
        is_header = r == 0
        for c in range(n_cols):
            if c >= len(row):
                continue
            cell_obj = row[c]
            if not isinstance(cell_obj, dict):
                continue
            content = cell_obj.get("content") or {}
            cell = table.cell(r, c)
            cell.margin_left = 0
            cell.margin_right = 0
            cell.margin_top = 0
            cell.margin_bottom = 0
            tf = cell.text_frame
            tf.word_wrap = False
            if is_header:
                if header_fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_fill
            # 表格边框（v1 tableStyles 四元组 [上,右,下,左]，逐 cell 应用）
            if border:
                _apply_cell_border(cell, border, theme_colors)
            # 单元格文本：富文本 <p> 分段 + <strong>/<em> 行内强调（还原粗体，
            # 防标签字面泄漏到 PPT）
            paragraphs = _parse_rich_text(str(content.get("text", ""))) or [("", False, False)]
            halign = content.get("align")
            align_val = None
            if isinstance(halign, (list, tuple)) and halign:
                align_val = _ALIGN_H_MAP.get(str(halign[0]).lower())
            for pi, runs in enumerate(paragraphs):
                p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                if align_val:
                    p.alignment = getattr(PP_ALIGN, align_val)
                for t, bold, italic in runs:
                    run = p.add_run()
                    run.text = t
                    run.font.size = Pt(int(font_size))
                    run.font.name = font_family
                    run.font.bold = bold
                    run.font.italic = italic
                    if is_header:
                        run.font.bold = header_bold
                        if header_color:
                            run.font.color.rgb = header_color
                    else:
                        color = _color_value(content.get("color"), theme_colors)
                        if color:
                            run.font.color.rgb = color
                        elif body_color:
                            run.font.color.rgb = body_color


# ---------------------------------------------------------------------------
# image 元素
# ---------------------------------------------------------------------------
def _render_image(slide, el, deck_dir):
    """图片：src（v1 要求绝对路径，相对路径按工程目录解析）+ fit 等比缩放。"""
    from pptx.util import Emu
    src = el.get("src")
    if not src:
        return
    if not os.path.isabs(src):
        src = os.path.join(deck_dir, src)
    if not os.path.exists(src):
        print(f"[convert warning] 图片不存在，跳过: {src}")
        return
    x, y, w, h = [int(v) for v in el["bounds"]]
    fit = (el.get("fit") or {}).get("mode", "fill")
    try:
        if fit in ("contain", "cover"):
            _add_picture_fit(slide, src, x, y, w, h, fit)
        else:
            slide.shapes.add_picture(src, Emu(x * EMU_PER_PT), Emu(y * EMU_PER_PT),
                                     Emu(w * EMU_PER_PT), Emu(h * EMU_PER_PT))
    except Exception as e:
        print(f"[convert warning] 图片插入失败，跳过: {e}")


def _add_picture_fit(slide, src, x, y, w, h, mode):
    """等比缩放图片到 bounds（contain 完整显示留白；cover 铺满并裁切）。

    用 PIL 读原始尺寸（主 venv 已装），PIL 不可用时降级为 fill 拉伸。
    """
    from pptx.util import Emu
    try:
        from PIL import Image
        with Image.open(src) as im:
            iw, ih = im.size
    except Exception:
        slide.shapes.add_picture(src, Emu(x * EMU_PER_PT), Emu(y * EMU_PER_PT),
                                 Emu(w * EMU_PER_PT), Emu(h * EMU_PER_PT))
        return
    if iw <= 0 or ih <= 0:
        slide.shapes.add_picture(src, Emu(x * EMU_PER_PT), Emu(y * EMU_PER_PT),
                                 Emu(w * EMU_PER_PT), Emu(h * EMU_PER_PT))
        return
    img_ratio = iw / ih
    if mode == "cover":
        # 图片更宽 -> 高度填满宽度溢出；否则宽度填满高度溢出
        disp_w = max(w, h * img_ratio)
        disp_h = max(h, w / img_ratio)
    else:  # contain
        disp_w = min(w, h * img_ratio)
        disp_h = min(h, w / img_ratio)
    dx = (w - disp_w) / 2
    dy = (h - disp_h) / 2
    slide.shapes.add_picture(
        src, Emu(int((x + dx) * EMU_PER_PT)), Emu(int((y + dy) * EMU_PER_PT)),
        Emu(int(disp_w * EMU_PER_PT)), Emu(int(disp_h * EMU_PER_PT)))


def _apply_cell_border(cell, border, theme_colors):
    """应用单元格边框。border = v1 四元组 [上,右,下,左]，每项 {style,width,color} 或 null。

    python-pptx 无 cell border API，操作 tcPr 的 a:lnT/lnR/lnB/lnL。
    """
    from pptx.oxml.ns import qn
    if not isinstance(border, (list, tuple)) or len(border) != 4:
        return
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag, spec in zip(("a:lnT", "a:lnR", "a:lnB", "a:lnL"), border, strict=False):
        if not isinstance(spec, dict):
            continue
        ln = tc_pr.find(qn(tag))
        if ln is None:
            ln = tc_pr.makeelement(qn(tag), {})
            tc_pr.append(ln)
        for child in list(ln):
            ln.remove(child)
        color = _color_value(spec.get("color"), theme_colors)
        if color:
            fill = ln.makeelement(qn("a:solidFill"), {})
            srgb = ln.makeelement(qn("a:srgbClr"), {"val": str(color)})
            fill.append(srgb)
            ln.append(fill)
        if spec.get("width"):
            ln.set("w", str(int(float(spec["width"]) * EMU_PER_PT)))
        if spec.get("style") == "dash":
            dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
            ln.append(dash)
