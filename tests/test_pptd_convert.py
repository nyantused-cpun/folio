# -*- coding: utf-8 -*-
"""pptd v1 方言 -> PPTX 转换器（_pptd_convert）单元测试。

覆盖：富文本解析、颜色解析、shapeName 映射、convert 产出结构（文本/形状/
连接线箭头/表格）、check 边界。全部走 tmp_path，不触真实 output/ 与 PowerPoint。
"""

import os


def _write_deck(tmp_path, main_text, pages):
    """写最小 pptd 工程，返回主文件绝对路径。"""
    root = tmp_path / "deck"
    (root / "pages").mkdir(parents=True)
    main = root / "deck.pptd"
    main.write_text(main_text, encoding="utf-8")
    for name, text in pages.items():
        (root / "pages" / name).write_text(text, encoding="utf-8")
    return str(main)


_OK_MAIN = "size: [1280, 720]\npages:\n- pages/01.page\n"


# ---------------------------------------------------------------------------
# _parse_rich_text
# ---------------------------------------------------------------------------
class TestParseRichText:
    def test_plain_text(self):
        from _pptd_convert import _parse_rich_text
        assert _parse_rich_text("hello") == [[("hello", False, False)]]

    def test_plain_text_trailing_newline_stripped(self):
        """v1 body 固定末尾带 \\n（pe.text 拼接），无 <p> 时应 strip 掉。"""
        from _pptd_convert import _parse_rich_text
        assert _parse_rich_text("结构化声明\n\n      ") == \
            [[("结构化声明", False, False)]]

    def test_bold_strong(self):
        from _pptd_convert import _parse_rich_text
        assert _parse_rich_text("<p><strong>标题</strong></p>") == \
            [[("标题", True, False)]]

    def test_mixed_runs(self):
        from _pptd_convert import _parse_rich_text
        result = _parse_rich_text("<strong>粗</strong>正常<em>斜</em>")
        # 三个 run：粗体、普通、斜体
        assert result == [[("粗", True, False), ("正常", False, False),
                           ("斜", False, True)]]

    def test_multi_paragraph(self):
        from _pptd_convert import _parse_rich_text
        result = _parse_rich_text("<p>第一行</p>\n<p>第二行</p>")
        assert len(result) == 2
        assert result[0] == [("第一行", False, False)]
        assert result[1] == [("第二行", False, False)]

    def test_empty(self):
        from _pptd_convert import _parse_rich_text
        assert _parse_rich_text("") == []


# ---------------------------------------------------------------------------
# _color_value
# ---------------------------------------------------------------------------
class TestColorValue:
    def test_hex(self):
        from _pptd_convert import _color_value
        from pptx.dml.color import RGBColor
        assert _color_value("#1E3A8A", {}) == RGBColor.from_string("1E3A8A")

    def test_token(self):
        from _pptd_convert import _color_value
        from pptx.dml.color import RGBColor
        assert _color_value("$primary", {"primary": "#2563EB"}) == \
            RGBColor.from_string("2563EB")

    def test_missing_token_returns_none(self):
        from _pptd_convert import _color_value
        assert _color_value("$nope", {}) is None

    def test_hex8_alpha_stripped(self):
        from _pptd_convert import _color_value
        from pptx.dml.color import RGBColor
        assert _color_value("#0A1540CC", {}) == RGBColor.from_string("0A1540")


# ---------------------------------------------------------------------------
# shapeName 映射覆盖
# ---------------------------------------------------------------------------
class TestShapeMap:
    def test_all_mapped_names_are_valid(self):
        """映射表每个成员名都必须是 MSO_SHAPE 的有效属性。"""
        from pptx.enum.shapes import MSO_SHAPE
        from _pptd_convert import _SHAPE_NAME_MAP, _CONNECTOR_MAP
        for member in _SHAPE_NAME_MAP.values():
            assert hasattr(MSO_SHAPE, member), f"MSO_SHAPE 缺 {member}"
        from pptx.enum.shapes import MSO_CONNECTOR
        for member in _CONNECTOR_MAP.values():
            assert hasattr(MSO_CONNECTOR, member), f"MSO_CONNECTOR 缺 {member}"


# ---------------------------------------------------------------------------
# convert 产出结构
# ---------------------------------------------------------------------------
class TestConvert:
    def _convert(self, tmp_path, page_text):
        main = _write_deck(tmp_path, _OK_MAIN, {"01.page": page_text})
        out = str(tmp_path / "out.pptx")
        from _pptd_convert import convert_pptd
        assert convert_pptd(main, out) is True
        assert os.path.exists(out)
        return out

    def test_text_and_shape(self, tmp_path):
        page = (
            "elements:\n"
            "- elementId: t1\n  elementType: text\n  bounds: [100, 100, 300, 60]\n"
            "  content: {text: '<p><strong>Hi</strong> world</p>', fontSize: 20}\n"
            "- elementId: s1\n  elementType: shape\n  bounds: [50, 50, 200, 100]\n"
            "  shapeName: roundRect\n  fill: {type: solid, color: '#1E3A8A'}\n")
        out = self._convert(tmp_path, page)
        from pptx import Presentation
        prs = Presentation(out)
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        # text + shape 两个元素
        assert len(shapes) == 2

    def test_connector_arrow(self, tmp_path):
        """连接线带箭头：tailEnd=triangle 落在终点。"""
        page = (
            "elements:\n"
            "- elementId: c1\n  elementType: shape\n  bounds: [100, 100, 200, 0]\n"
            "  shapeName: straightConnector1\n"
            "  border: {style: solid, width: 1.5, color: '#4A6BB8'}\n"
            "  flip: [false, false]\n"
            "  arrow: [none, arrow]\n")
        out = self._convert(tmp_path, page)
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(out)
        slide = prs.slides[0]
        conns = [s for s in slide.shapes if str(s.shape_type).startswith(("LINE", "CONNECTOR"))]
        assert len(conns) == 1
        ln = conns[0].line._get_or_add_ln()
        tail = ln.find(qn("a:tailEnd"))
        assert tail is not None and tail.get("type") == "triangle"

    def test_connector_stealth_arrow_uses_valid_enum(self, tmp_path):
        """stealth 箭头必须产出合法 ST_LineEndType 值 'stealth'（非 'stealthArrow'）。

        回归：2026-08-12 蓝海集团真实工程实测，'stealthArrow' 非法导致
        PowerPoint 拒绝打开 pptx。
        """
        page = (
            "elements:\n"
            "- elementId: c1\n  elementType: shape\n  bounds: [100, 100, 200, 0]\n"
            "  shapeName: straightConnector1\n"
            "  border: {style: solid, width: 1.5, color: '#4A6BB8'}\n"
            "  flip: [false, false]\n"
            "  arrow: [stealth, arrow]\n")
        out = self._convert(tmp_path, page)
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(out)
        slide = prs.slides[0]
        conn = [s for s in slide.shapes if str(s.shape_type).startswith(("LINE", "CONNECTOR"))][0]
        ln = conn.line._get_or_add_ln()
        head = ln.find(qn("a:headEnd"))
        tail = ln.find(qn("a:tailEnd"))
        assert head is not None and head.get("type") == "stealth"
        assert tail is not None and tail.get("type") == "triangle"

    def test_table(self, tmp_path):
        page = (
            "elements:\n"
            "- elementId: tb1\n  elementType: table\n  bounds: [80, 128, 1120, 300]\n"
            "  columnWidths: [0.3, 0.7]\n  rowHeights: [0.5, 0.5]\n"
            "  rows:\n"
            "  - - content: {text: '序号'}\n    - content: {text: '内容'}\n"
            "  - - content: {text: '1'}\n    - content: {text: 'A'}\n")
        out = self._convert(tmp_path, page)
        from pptx import Presentation
        prs = Presentation(out)
        slide = prs.slides[0]
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        assert len(tables[0].table.rows) == 2
        assert len(tables[0].table.columns) == 2

    def test_unknown_shape_falls_back_rect(self, tmp_path, capsys):
        page = (
            "elements:\n"
            "- elementId: s1\n  elementType: shape\n  bounds: [50, 50, 100, 100]\n"
            "  shapeName: 不存在的形状\n")
        out = self._convert(tmp_path, page)
        assert os.path.exists(out)
        assert "未知 shapeName" in capsys.readouterr().out

    def test_shape_flip_sets_xfrm_attr(self, tmp_path):
        """shape flip 应落到 a:xfrm 的 flipH/flipV（timeline 红三角 flipV）。"""
        page = (
            "elements:\n"
            "- elementId: s1\n  elementType: shape\n  bounds: [50, 50, 100, 100]\n"
            "  shapeName: triangle\n"
            "  fill: {type: solid, color: '#FF0000'}\n"
            "  flip: [false, true]\n")
        out = self._convert(tmp_path, page)
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(out)
        shape = [s for s in prs.slides[0].shapes
                 if str(s.shape_type).startswith("AUTO_SHAPE")][0]
        xfrm = shape._element.spPr.find(qn("a:xfrm"))
        assert xfrm is not None and xfrm.get("flipV") == "1"

    def test_table_border_applied(self, tmp_path):
        """v1 tableStyles border 应落到 cell 的 lnB（底部分隔线）。"""
        main = (
            "size: [1280, 720]\n"
            "theme:\n"
            "  colors:\n    hairline: '#D8E2EE'\n"
            "  tableStyles:\n"
            "    default:\n"
            "      border:\n      - null\n      - null\n"
            "      - style: solid\n        width: 0.75\n        color: '$hairline'\n"
            "      - null\n"
            "pages:\n- pages/01.page\n")
        page = (
            "elements:\n"
            "- elementId: tb1\n  elementType: table\n  bounds: [80, 128, 1120, 300]\n"
            "  columnWidths: [0.3, 0.7]\n  rowHeights: [0.5, 0.5]\n"
            "  style: '$default'\n"
            "  rows:\n"
            "  - - content: {text: '序号'}\n    - content: {text: '内容'}\n"
            "  - - content: {text: '1'}\n    - content: {text: 'A'}\n")
        main_path = _write_deck(tmp_path, main, {"01.page": page})
        out = str(tmp_path / "out.pptx")
        from _pptd_convert import convert_pptd
        assert convert_pptd(main_path, out) is True
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(out)
        table = [s for s in prs.slides[0].shapes if s.has_table][0].table
        cell = table.cell(0, 0)
        lnB = cell._tc.get_or_add_tcPr().find(qn("a:lnB"))
        assert lnB is not None

    def test_missing_page_skipped_with_warning(self, tmp_path, capsys):
        main = _write_deck(tmp_path, "size: [1280, 720]\npages:\n- pages/不存在.page\n", {})
        out = str(tmp_path / "out.pptx")
        from _pptd_convert import convert_pptd
        assert convert_pptd(main, out) is True
        assert "页面文件不存在" in capsys.readouterr().out


# ---------------------------------------------------------------------------
# check 边界
# ---------------------------------------------------------------------------
class TestCheck:
    def test_valid_deck(self, tmp_path):
        from _pptd_convert import check_pptd
        main = _write_deck(
            tmp_path, _OK_MAIN,
            {"01.page": "elements:\n- elementId: s1\n  elementType: shape\n"
                        "  bounds: [0, 0, 10, 10]\n  shapeName: rect\n"})
        ok, warn = check_pptd(main)
        assert ok is True
        assert warn == {}

    def test_bad_size(self, tmp_path, capsys):
        from _pptd_convert import check_pptd
        main = _write_deck(tmp_path, "size: 1280\npages: []\n", {})
        assert check_pptd(main)[0] is False
        assert "size 必须是" in capsys.readouterr().out

    def test_bad_bounds(self, tmp_path, capsys):
        from _pptd_convert import check_pptd
        main = _write_deck(
            tmp_path, _OK_MAIN,
            {"01.page": "elements:\n- elementId: s1\n  elementType: shape\n"
                        "  bounds: [0, 0]\n  shapeName: rect\n"})
        assert check_pptd(main)[0] is False
        assert "bounds 必须是" in capsys.readouterr().out
